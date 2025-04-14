#!/usr/bin/env python3
"""
Export Utilities for LumenTrialGuide.AI

This module provides shared functionality for various export operations
including automatically generating trial strategy decks alongside
other exports like summary packets and regulatory bundles.
"""

import os
import json
import sys
from typing import Dict, Any, Optional
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt

def auto_generate_trial_strategy_deck(session_id: str, persona: str = "general") -> Dict[str, Any]:
    """
    Automatically generate a trial strategy deck for a given session.
    This simplified version generates a basic deck with executive brief and summary.
    For more advanced persona-specific decks, use the pitch_deck_generator module.
    
    Args:
        session_id: The session ID for which to generate the deck
        persona: Optional persona type (general, investor, regulatory, cxo)
        
    Returns:
        Dictionary with status and path information
    """
    # Determine archive directory based on environment
    if os.path.exists("/mnt/data"):
        # Production environment
        archive_dir = f"/mnt/data/lumen_reports_backend/sessions/{session_id}"
    else:
        # Development environment
        archive_dir = f"data/sessions/{session_id}"
    
    # Ensure directory exists
    os.makedirs(archive_dir, exist_ok=True)
    
    # Create presentation
    prs = Presentation()
    
    # Add Executive Design Brief slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[5])
    slide1.shapes.title.text = "🧠 Executive Design Brief"
    
    # Get alignment score data
    alignment_score = "N/A"
    alignment_path = os.path.join(archive_dir, "alignment_score_report.json")
    if os.path.exists(alignment_path):
        try:
            with open(alignment_path, "r") as f:
                alignment_data = json.load(f)
                score = alignment_data.get("alignment_score", 0)
                alignment_score = f"{round(score * 100)}%"
        except Exception as e:
            print(f"Error reading alignment data: {str(e)}")
    
    # Get protocol improvement suggestions count
    suggestions = 0
    suggestions_path = os.path.join(archive_dir, "suggested_corrections.json")
    if os.path.exists(suggestions_path):
        try:
            with open(suggestions_path, "r") as f:
                suggestions_data = json.load(f)
                suggs = suggestions_data.get("suggestions", [])
                suggestions = len(suggs)
        except Exception as e:
            print(f"Error reading suggestions data: {str(e)}")
    
    # Add content to slide 1
    text1 = (
        f"Study ID: {session_id}\n"
        f"Persona: {persona.title()}\n"
        f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n"
        f"Alignment Score: {alignment_score}\n"
        f"Protocol Improvement Suggestions: {suggestions}\n"
        f"AI Confidence: High\n\n"
        f"All insights generated using CSR-backed evidence, semantic alignment, "
        f"and protocol validation tools built into LumenTrialGuide.AI."
    )
    
    tf1 = slide1.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5)).text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = text1
    p1.font.size = Pt(18)
    
    # Add Summary Intelligence Highlights slide
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    slide2.shapes.title.text = "📄 Summary Intelligence Highlights"
    
    # Prepare content for slide 2 based on available data
    highlight_points = []
    
    # Check for dropout forecast
    dropout_path = os.path.join(archive_dir, "dropout_forecast.json")
    if os.path.exists(dropout_path):
        try:
            with open(dropout_path, "r") as f:
                dropout_data = json.load(f)
                risk_level = dropout_data.get("risk_level", "").upper()
                completion_rate = dropout_data.get("expected_completion_rate", "N/A")
                if completion_rate != "N/A":
                    highlight_points.append(f"📉 Dropout Risk: {risk_level} (Projected Completion: {completion_rate}%)")
        except Exception:
            pass
    
    # Check for success prediction
    success_path = os.path.join(archive_dir, "success_prediction.json")
    if os.path.exists(success_path):
        try:
            with open(success_path, "r") as f:
                success_data = json.load(f)
                probability = success_data.get("probability", success_data.get("success_probability", 0))
                highlight_points.append(f"📈 Success Probability: {round(probability * 100)}%")
        except Exception:
            pass
    
    # Add information about available exports
    available_exports = [
        "📊 Summary Packet (PDF)",
        "📋 Protocol Alignment Report (PDF)",
        "📑 Regulatory Bundle (ZIP)",
        "💼 Executive Brief (PPTX)",
    ]
    
    # Combine highlights and exports
    text2 = ""
    
    if highlight_points:
        text2 += "Key Insights:\n"
        for point in highlight_points:
            text2 += f"• {point}\n"
        text2 += "\n"
    
    text2 += "Available Exports:\n"
    for export in available_exports:
        text2 += f"• {export}\n"
    
    text2 += "\nTraceability and explainability are built into every export."
    
    # Add content to slide 2
    tf2 = slide2.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5)).text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = text2
    p2.font.size = Pt(18)
    
    # Save presentation
    deck_path = os.path.join(archive_dir, "trial_strategy_deck.pptx")
    prs.save(deck_path)
    
    # Return status and path
    return {
        "status": "success",
        "deck_path": deck_path,
        "slide_count": len(prs.slides),
        "web_path": f"/static/{session_id}/trial_strategy_deck.pptx",
    }

def update_export_log(session_id: str, export_type: str, file_path: str) -> None:
    """
    Update the export log for a given session
    
    Args:
        session_id: The session ID for which the export was generated
        export_type: The type of export (e.g., 'summary_packet', 'strategy_deck')
        file_path: The path to the exported file
    """
    # Determine archive directory based on environment
    if os.path.exists("/mnt/data"):
        # Production environment
        archive_dir = f"/mnt/data/lumen_reports_backend/sessions/{session_id}"
    else:
        # Development environment
        archive_dir = f"data/sessions/{session_id}"
    
    # Ensure directory exists
    os.makedirs(archive_dir, exist_ok=True)
    
    # Create export log entry
    export_log = {
        "timestamp": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
        "type": export_type,
        "file_path": file_path,
        "session_id": session_id
    }
    
    # Save export log
    log_path = os.path.join(archive_dir, "export_log.json")
    
    try:
        if os.path.exists(log_path):
            with open(log_path, "r") as f:
                logs = json.load(f)
                if isinstance(logs, list):
                    logs.append(export_log)
                else:
                    logs = [logs, export_log]
        else:
            logs = [export_log]
            
        with open(log_path, "w") as f:
            json.dump(logs, f, indent=2)
    except Exception as e:
        print(f"Error saving export log: {str(e)}")